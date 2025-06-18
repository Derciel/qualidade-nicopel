import streamlit as st
import streamlit_authenticator as stauth
import yaml
from yaml.loader import SafeLoader
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import gspread
from google.oauth2.service_account import Credentials
import io
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import requests
from PIL import Image
import plotly.graph_objects as go

st.set_page_config(page_title="Dashboard de Não Conformidades", page_icon="📊", layout="wide")

# --- Autenticação com streamlit_authenticator ---
config = {
    "credentials": {
        "usernames": {
            "ti": {
                "email": "ti@nicopel.com.br",
                "name": "TI",
                "password": "$2b$12$XXXXXXXXXXXXXXXXXXXXXXXXXXXXXX"
            },
            "qualidade": {
                "email": "qualidade@nicopel.com.br",
                "name": "Qualidade",
                "password": "$2b$12$YYYYYYYYYYYYYYYYYYYYYYYYYYYYYY"
            }
        }
    },
    "cookie": {
        "name": "dashboard_cookie",
        "key": "random_signature_key",
        "expiry_days": 1
    },
    "preauthorized": {
        "emails": ["ti@nicopel.com.br", "qualidade@nicopel.com.br"]
    }
}

authenticator = stauth.Authenticate(
    config["credentials"],
    config["cookie"]["name"],
    config["cookie"]["key"],
    config["cookie"]["expiry_days"]
)

if authentication_status is False:
    st.error("Usuário ou senha incorretos.")
elif authentication_status is None:
    st.warning("Por favor, insira suas credenciais.")
elif authentication_status:
    authenticator.logout("Sair", "sidebar")
    st.sidebar.success(f"Logado como {name}")

    # --- Carregamento de Dados Google Sheets ---
    GOOGLE_SHEETS_CREDENTIALS = st.secrets["gcp_service_account"]
    NOME_PLANILHA = '16KWu85cbnA6wxY8pjEbyAAqBUuxE9iUmOCmdAhKSF9Y'
    NOME_ABA = 'Form'

    @st.cache_data(ttl=300)
    def load_data_from_gsheets():
        try:
            scope = [
                'https://www.googleapis.com/auth/spreadsheets',
                'https://www.googleapis.com/auth/drive.readonly'
            ]
            creds = Credentials.from_service_account_info(GOOGLE_SHEETS_CREDENTIALS, scopes=scope)
            gc = gspread.authorize(creds)
            planilha = gc.open_by_key(NOME_PLANILHA)
            aba = planilha.worksheet(NOME_ABA)
            dados = aba.get_all_records()
            if not dados:
                st.warning("A aba da planilha está vazia.")
                return pd.DataFrame()
            df = pd.DataFrame(dados)
            df.rename(columns={"CLASSIFICAÇAO NC": "CLASSIFICAÇÃO NC"}, inplace=True)
            df.columns = [col.strip() for col in df.columns]

            df['DATA DA NAO CONFORMIDADE'] = pd.to_datetime(df['DATA DA NAO CONFORMIDADE'], errors='coerce', dayfirst=True)
            df['DATA DE ENCERRAMENTO'] = pd.to_datetime(df['DATA DE ENCERRAMENTO'], errors='coerce', dayfirst=True)
            df.dropna(subset=['DATA DA NAO CONFORMIDADE'], inplace=True)
            df['STATUS'] = df['DATA DE ENCERRAMENTO'].apply(lambda x: 'Resolvida' if pd.notna(x) else 'Pendente')

            cols_to_str = ['CLIENTE (Caso tenha)', 'DEPARTAMENTO RESPONSÁVEL', 'SETOR DO RESPONSÁVEL', 'CLASSIFICAÇÃO NC', 'AVALIAÇÃO DA EFICÁCIA']
            for col in cols_to_str:
                if col in df.columns:
                    df[col] = df[col].astype(str)
            return df
        except Exception as e:
            st.error(f"Ocorreu um erro ao carregar os dados: {e}")
            return pd.DataFrame()

    st.title("📊 Dashboard de Análise de Não Conformidades")
    df = load_data_from_gsheets()
    st.dataframe(df)


def download_image_from_url(url):
    try:
        response = requests.get(url)
        img = Image.open(io.BytesIO(response.content))
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        img.save(temp_file.name)
        return temp_file.name
    except:
        return None

# ALTERADO: Função de PPTX agora aceita o dicionário de cores

def create_powerpoint_presentation(df, logo_url, cores_departamentos):
    prs = Presentation()
    # Slide 1: Título (sem alterações)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Relatório de Não Conformidades"
    slide.placeholders[1].text = f"Gerado em: {pd.Timestamp.now().strftime('%d/%m/%Y %H:%M')}"
    logo_path = download_image_from_url(logo_url)
    if logo_path:
        slide.shapes.add_picture(logo_path, Inches(8), Inches(0.2), width=Inches(1.5))

    # --- NOVO: Slide 2 com os KPIs como Gráficos de Rosca (Donut) ---
    kpi_slide = prs.slides.add_slide(prs.slide_layouts[5])
    kpi_slide.shapes.title.text = "Principais Indicadores (KPIs)"

    # Calcular os valores e porcentagens dos KPIs
    total_ncs = df.shape[0] if not df.empty else 1
    ncs_pendentes = df[df['STATUS'] == 'Pendente'].shape[0]
    ncs_resolvidas = df[df['STATUS'] == 'Resolvida'].shape[0]
    
    perc_pendentes = (ncs_pendentes / total_ncs) * 100
    perc_resolvidas = (ncs_resolvidas / total_ncs) * 100

    kpis = [
        {'label': 'NCs Pendentes', 'value': ncs_pendentes, 'percent': perc_pendentes, 'color': RGBColor(228, 87, 87), 'total': total_ncs},
        {'label': 'NCs Resolvidas', 'value': ncs_resolvidas, 'percent': perc_resolvidas, 'color': RGBColor(87, 163, 105), 'total': total_ncs}
    ]

    # Posições iniciais para os gráficos
    left = Inches(1.5)
    top = Inches(1.5)
    size = Inches(3.0) # Tamanho do diâmetro
    gap = Inches(1.0)

    # Criar um "velocímetro" para cada KPI
    for kpi in kpis:
        # 1. Círculo de fundo (cinza claro)
        kpi_slide.shapes.add_shape(MSO_SHAPE.DONUT, left, top, size, size)
        
        # 2. Arco de progresso (colorido)
        # O ajuste `adj2` controla a "abertura" do arco. 21600000 = 360 graus.
        # O ajuste `adj1` controla o tamanho do "buraco" do donut.
        arc = kpi_slide.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, left, top, size, size)
        arc.adjustments[0] = -5400000  # Começa no topo (12h)
        arc.adjustments[1] = int((kpi['percent'] / 100) * 21600000 - 5400000) # Ângulo final
        arc.line.fill.background()
        arc.fill.solid()
        arc.fill.fore_color.rgb = kpi['color']
        
        # 3. Texto no centro com a porcentagem
        textbox_percent = kpi_slide.shapes.add_textbox(left, top, size, size)
        tf_percent = textbox_percent.text_frame
        tf_percent.clear()
        tf_percent.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_percent = tf_percent.paragraphs[0]
        p_percent.text = f"{kpi['percent']:.0f}%"
        p_percent.font.name = 'Roboto Slab'
        p_percent.font.size = Pt(40)
        p_percent.font.bold = True
        p_percent.alignment = PP_ALIGN.CENTER
        
        # 4. Texto abaixo com o rótulo e os valores
        textbox_label = kpi_slide.shapes.add_textbox(left, top + size - Inches(0.2), size, Inches(1))
        tf_label = textbox_label.text_frame
        tf_label.clear()
        p_label = tf_label.paragraphs[0]
        p_label.text = f"{kpi['label']}\n({kpi['value']} de {kpi['total']})"
        p_label.font.name = 'Roboto Slab'
        p_label.font.size = Pt(14)
        p_label.alignment = PP_ALIGN.CENTER

        left += size + gap

    fig, ax = plt.subplots()
    df['CLASSIFICAÇÃO NC'].value_counts().plot.pie(autopct='%1.1f%%', ax=ax)
    ax.set_ylabel('')
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png")
    buf.seek(0)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Classificação das NCs"
    slide.shapes.add_picture(buf, Inches(2), Inches(1.5), width=Inches(6))

    # Slide 4 - Gráfico Departamentos (COM CORES PERSONALIZADAS)
    fig_depto, ax_depto = plt.subplots(figsize=(10, 6))
    depto_data = df.groupby('DEPARTAMENTO RESPONSÁVEL').size().reset_index(name='Quantidade')
    sns.barplot(data=depto_data, x='DEPARTAMENTO RESPONSÁVEL', y='Quantidade', palette=cores_departamentos, ax=ax_depto)
    ax_depto.set_title("NCs por Departamento", fontsize=14, weight='bold')
    ax_depto.set_xlabel("Departamento", fontsize=12)
    ax_depto.set_ylabel("Quantidade", fontsize=12)
    plt.xticks(rotation=45, ha='right')
    plt.tight_layout()
    buf_depto = io.BytesIO()
    fig_depto.savefig(buf_depto, format="png")
    buf_depto.seek(0)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "NCs por Departamento"
    slide.shapes.add_picture(buf_depto, Inches(0.5), Inches(1.5), width=Inches(9))

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    return pptx_io.getvalue()

# NOVO: Função para criar KPIs de medidor (gauge)
def create_gauge_chart(value, title, max_value, color):
    fig = go.Figure(go.Indicator(
        mode="gauge+number",
        value=value,
        title={'text': title, 'font': {'size': 20, 'family': 'Roboto Slab'}},
        gauge={
            'axis': {'range': [0, max_value], 'tickwidth': 1, 'tickcolor': "darkblue"},
            'bar': {'color': color},
            'steps': [
                {'range': [0, max_value * 0.5], 'color': 'lightgray'},
                {'range': [max_value * 0.5, max_value], 'color': 'darkgray'}
            ],
        }
    ))
    fig.update_layout(
        paper_bgcolor='rgba(0,0,0,0)',
        font={'color': "white" if tema_selecionado == "Escuro" else "black", 'family': "Roboto Slab"}
    )
    return fig

# --- INTERFACE PRINCIPAL ---
st.title("📊 Dashboard de Análise de Não Conformidades")
df = load_data_from_gsheets()
if df.empty:
    st.stop()

# --- FILTROS NA SIDEBAR ---
st.sidebar.header("Filtros Interativos")
min_date = df['DATA DA NAO CONFORMIDADE'].min().date()
max_date = df['DATA DA NAO CONFORMIDADE'].max().date()
date_range = st.sidebar.date_input("Período da Não Conformidade:", value=(min_date, max_date), min_value=min_date, max_value=max_date)
selected_classificacao = st.sidebar.multiselect("Classificação NC:", options=df['CLASSIFICAÇÃO NC'].unique(), default=df['CLASSIFICAÇÃO NC'].unique())
selected_departments = st.sidebar.multiselect("Departamento Responsável:", options=df['DEPARTAMENTO RESPONSÁVEL'].unique(), default=df['DEPARTAMENTO RESPONSÁVEL'].unique())
selected_status = st.sidebar.multiselect("Status:", options=df['STATUS'].unique(), default=df['STATUS'].unique())

start_date, end_date = date_range
df_filtered = df[
    (df['DATA DA NAO CONFORMIDADE'].dt.date >= start_date) &
    (df['DATA DA NAO CONFORMIDADE'].dt.date <= end_date) &
    (df['CLASSIFICAÇÃO NC'].isin(selected_classificacao)) &
    (df['DEPARTAMENTO RESPONSÁVEL'].isin(selected_departments)) &
    (df['STATUS'].isin(selected_status))
]

# NOVO: Seção de seleção de cores na sidebar
st.sidebar.markdown("---")
st.sidebar.subheader("Cores por Departamento")
# Usamos o DF original para ter todos os departamentos sempre disponíveis para colorir
departamentos_unicos = df['DEPARTAMENTO RESPONSÁVEL'].unique()
cores_setores = {}
for depto in departamentos_unicos:
    # A chave do color_picker precisa ser única
    cor_padrao = "#1f77b4" # Azul padrão
    cores_setores[depto] = st.sidebar.color_picker(f"Cor para {depto}", value=cor_padrao, key=f"color_{depto}")

# BOTÃO DE EXPORTAÇÃO PARA PPTX
st.sidebar.markdown("---")
st.sidebar.header("Exportar Relatório")
if not df_filtered.empty:
    # ALTERADO: Passando o dicionário de cores para a função do PowerPoint
    pptx_bytes = create_powerpoint_presentation(df_filtered, LOGO_URL, cores_setores)
    st.sidebar.download_button(
        label="Exportar para PowerPoint (.pptx)",
        data=pptx_bytes,
        file_name="Relatorio_NaoConformidades.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
else:
    st.sidebar.info("Não há dados para exportar com os filtros atuais.")

# ALTERADO: KPIs agora são medidores visuais
st.markdown("---")
total_ncs = df_filtered.shape[0]
ncs_pendentes = df_filtered[df_filtered['STATUS'] == 'Pendente'].shape[0]
df_resolvidas = df_filtered[df_filtered['STATUS'] == 'Resolvida']
taxa_resolucao = (len(df_resolvidas) / total_ncs * 100) if total_ncs > 0 else 0

col1, col2, col3 = st.columns(3)
with col1:
    st.plotly_chart(create_gauge_chart(total_ncs, "Total de NCs", max_value=len(df), color="#0083B8"), use_container_width=True)
with col2:
    st.plotly_chart(create_gauge_chart(ncs_pendentes, "NCs Pendentes", max_value=total_ncs, color="#E45757"), use_container_width=True)
with col3:
    st.plotly_chart(create_gauge_chart(taxa_resolucao, "% NCs Resolvidas", max_value=100, color="#57A369"), use_container_width=True)


st.markdown("---")

# Gráficos de Pizza
col_graf1, col_graf2 = st.columns(2)
with col_graf1:
    st.subheader("Distribuição por Classificação")
    fig, ax = plt.subplots()
    df_filtered['CLASSIFICAÇÃO NC'].value_counts().plot.pie(ax=ax, autopct='%1.1f%%', startangle=90)
    ax.set_ylabel('')
    fig.patch.set_alpha(0.0) # Fundo transparente
    ax.patch.set_alpha(0.0)
    st.pyplot(fig)
with col_graf2:
    st.subheader("Avaliação da Eficácia das Ações")
    df_eficacia = df_filtered[df_filtered['AVALIAÇÃO DA EFICÁCIA'].replace('', pd.NA).notna()]
    if not df_eficacia.empty:
        fig, ax = plt.subplots()
        df_eficacia['AVALIAÇÃO DA EFICÁCIA'].value_counts().plot.pie(ax=ax, autopct='%1.1f%%', startangle=90)
        ax.set_ylabel('')
        fig.patch.set_alpha(0.0)
        ax.patch.set_alpha(0.0)
        st.pyplot(fig)
    else:
        st.write("Nenhum dado de eficácia para exibir.")

st.markdown("---")
st.subheader("Não Conformidades por Departamento")
if not df_filtered.empty:
    df_depto = df_filtered.groupby('DEPARTAMENTO RESPONSÁVEL').size().reset_index(name='Quantidade')
    
    fig, ax = plt.subplots(figsize=(10, 6))
    
    # ALTERADO: Gráfico de barras principal agora usa as cores personalizadas
    sns.barplot(
        data=df_depto,
        x='DEPARTAMENTO RESPONSÁVEL',
        y='Quantidade',
        palette=cores_setores, # Usando o dicionário de cores
        ax=ax
    )
    
    ax.set_title("Não Conformidades por Departamento", fontsize=14, weight='bold')
    ax.set_xlabel("Departamento", fontsize=12)
    ax.set_ylabel("Quantidade", fontsize=12)
    ax.set_xticklabels(ax.get_xticklabels(), rotation=45, ha='right')
    sns.despine()
    fig.patch.set_alpha(0.0)
    ax.patch.set_alpha(0.0)
    
    st.pyplot(fig)
else:
    st.write("Nenhum dado para exibir com os filtros atuais.")


st.subheader("Dados Detalhados (Filtrados)")
st.dataframe(df_filtered)